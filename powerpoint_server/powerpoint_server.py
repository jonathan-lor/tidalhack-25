import os
from typing import Any, List, Dict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from mcp.server.fastmcp import FastMCP

mcp = FastMCP("powerpoint_server")

@mcp.tool()
def create_powerpoint(topic: str, slides: List[Dict[str, Any]], output_file: str,
                      save_directory: str = None, title_slide: bool = True,
                      template_file: str = None) -> Dict[str, Any]:
    """
    Creates a styled PowerPoint presentation using existing slide placeholders.
    
    Parameters:
      topic (str): The main title of the presentation.
      slides (List[Dict[str, Any]]): A list of dictionaries defining each slide.
          Each slide dict can include:
            - "layout": (int) Slide layout index (default is 1).
            - "title": (str) Title text for the slide.
            - "content": (str or list) Text content. If string with newlines,
                         each line becomes a paragraph.
            - "bullets": (list) A list of bullet point strings.
            - "background_color": (str) Hex code (e.g., "#FF5733") for slide background.
      output_file (str): File name for the saved presentation.
      save_directory (str, optional): Directory to save the file.
      title_slide (bool): Whether to include a title slide.
      template_file (str, optional): Path to a PowerPoint template file (to use design themes).
      
    Returns:
      dict: Status message with success or error.
    """
    save_directory = "C:/Users/doria/Documents/VS/"

    try:
        # Use template if provided; otherwise, start with a blank presentation.
        prs = Presentation(template_file) if template_file else Presentation()

        # --- Title Slide ---
        if title_slide:
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = topic
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = "Generated by Gemini model"
            for para in slide.shapes.title.text_frame.paragraphs:
                para.font.size = Pt(40)
                para.alignment = PP_ALIGN.CENTER

        # --- Process Additional Slides ---
        for slide_info in slides:
            layout_idx = slide_info.get("layout", 1)
            slide_layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set slide title using the title placeholder if available.
            if "title" in slide_info:
                if slide.shapes.title:
                    slide.shapes.title.text = slide_info["title"]
                    for para in slide.shapes.title.text_frame.paragraphs:
                        para.font.size = Pt(32)
                else:
                    left, top, width, height = Inches(0.5), Inches(0.5), Inches(9), Inches(1)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    txBox.text = slide_info["title"]
                    for para in txBox.text_frame.paragraphs:
                        para.font.size = Pt(32)

            # --- If slide has bullets, add both content (if any) and bullet points ---
            if "bullets" in slide_info and isinstance(slide_info["bullets"], list):
                # Try to use the existing content placeholder.
                if len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]
                else:
                    content_placeholder = None

                if content_placeholder:
                    tf = content_placeholder.text_frame
                    tf.clear()
                    # If there's additional content, add it as the first paragraph.
                    if "content" in slide_info:
                        # Support both string and list types for content.
                        if isinstance(slide_info["content"], list):
                            content_text = " ".join(slide_info["content"])
                        elif isinstance(slide_info["content"], str):
                            content_text = slide_info["content"]
                        else:
                            content_text = ""
                        p = tf.add_paragraph()
                        p.text = content_text
                        p.font.size = Pt(24)
                        p.level = 0  # No bullet for content.
                        # Optional: Add an empty paragraph for spacing.
                        blank = tf.add_paragraph()
                        blank.text = ""
                    # Now add bullet points.
                    for bullet in slide_info["bullets"]:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 1  # Bullet level.
                        p.font.size = Pt(24)
                else:
                    # Fallback: create a new textbox if no placeholder exists.
                    left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(4)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    # Add content first if available.
                    if "content" in slide_info:
                        if isinstance(slide_info["content"], list):
                            p = tf.add_paragraph()
                            p.text = " ".join(slide_info["content"])
                        elif isinstance(slide_info["content"], str):
                            p = tf.add_paragraph()
                            p.text = slide_info["content"]
                        p.font.size = Pt(24)
                        p.level = 0
                        tf.add_paragraph()  # blank paragraph for spacing
                    # Now add bullet points.
                    for bullet in slide_info["bullets"]:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.level = 1
                        p.font.size = Pt(24)
            
            # --- If slide has only content (no bullets) ---
            elif "content" in slide_info:
                if len(slide.placeholders) > 1:
                    content_placeholder = slide.placeholders[1]
                else:
                    content_placeholder = None

                if content_placeholder:
                    tf = content_placeholder.text_frame
                    tf.clear()
                    if isinstance(slide_info["content"], list):
                        if slide_info["content"]:
                            p = tf.add_paragraph()
                            p.text = slide_info["content"][0]
                            p.font.size = Pt(24)
                        for line in slide_info["content"][1:]:
                            p = tf.add_paragraph()
                            p.text = line
                            p.font.size = Pt(24)
                    elif isinstance(slide_info["content"], str):
                        paragraphs = slide_info["content"].split("\n")
                        if paragraphs:
                            p = tf.add_paragraph()
                            p.text = paragraphs[0]
                            p.font.size = Pt(24)
                        for par in paragraphs[1:]:
                            p = tf.add_paragraph()
                            p.text = par
                            p.font.size = Pt(24)
                else:
                    # Fallback: add a textbox if no content placeholder exists.
                    left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(4)
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    if isinstance(slide_info["content"], list):
                        if slide_info["content"]:
                            p = tf.add_paragraph()
                            p.text = slide_info["content"][0]
                            p.font.size = Pt(24)
                        for line in slide_info["content"][1:]:
                            p = tf.add_paragraph()
                            p.text = line
                            p.font.size = Pt(24)
                    elif isinstance(slide_info["content"], str):
                        paragraphs = slide_info["content"].split("\n")
                        if paragraphs:
                            p = tf.add_paragraph()
                            p.text = paragraphs[0]
                            p.font.size = Pt(24)
                        for par in paragraphs[1:]:
                            p = tf.add_paragraph()
                            p.text = par
                            p.font.size = Pt(24)
            
            # --- Background color if provided ---
            if "background_color" in slide_info:
                hex_color = slide_info["background_color"].lstrip('#')
                if len(hex_color) == 6:
                    r, g, b = int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16)
                    slide.background.fill.solid()
                    slide.background.fill.fore_color.rgb = RGBColor(r, g, b)
        
        # Determine final save path.
        output_path = os.path.join(save_directory, output_file)
        prs.save(output_path)
        return {"status": "success", "message": f"PowerPoint file saved to {output_path}"}
    
    except Exception as e:
        return {"error": f"Failed to create PowerPoint presentation: {str(e)}"}

@mcp.tool()
def add_image_slide(presentation_file: str, image_path: str, title: str = None,
                    output_file: str = None, save_directory: str = None) -> Dict[str, Any]:
    """
    Adds an image slide to an existing presentation using an existing title placeholder if available.
    """
    save_directory = "C:/Users/doria/Documents/VS/"

    try:
        prs = Presentation(presentation_file)
        slide_layout = prs.slide_layouts[5]  # Blank layout.
        slide = prs.slides.add_slide(slide_layout)
        
        if title:
            title_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:  # Title placeholder type.
                    title_placeholder = shape
                    break
            if title_placeholder:
                title_placeholder.text = title
                for para in title_placeholder.text_frame.paragraphs:
                    para.font.size = Pt(32)
                    para.alignment = PP_ALIGN.CENTER
            else:
                left, top, width, height = Inches(1), Inches(0.5), Inches(8), Inches(1)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                txBox.text = title
                for para in txBox.text_frame.paragraphs:
                    para.font.size = Pt(32)
                    para.alignment = PP_ALIGN.CENTER
        
        left, top = Inches(1), Inches(2)
        slide.shapes.add_picture(image_path, left, top)
        
        if output_file:
            save_path = os.path.join(save_directory, output_file)
        else:
            save_path = presentation_file
        
        prs.save(save_path)
        return {"status": "success", "message": f"Image slide added and presentation saved to {save_path}"}
    
    except Exception as e:
        return {"error": f"Failed to add image slide: {str(e)}"}
    
if __name__ == "__main__":
    mcp.run()
